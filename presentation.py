from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_textbox(slide, text, left, top, width, height, font_size=12, is_special=False, is_header=False):
    """Adiciona caixa de texto padronizada com Arial"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), 
        Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if is_special else PP_ALIGN.LEFT
    
    if is_special:
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200, 200, 200)
    
    run = p.add_run()
    run.text = text
    run.font.name = 'Arial'  # Fonte Arial para todas as caixas
    run.font.size = Pt(font_size)
    run.font.bold = True if is_header else False  # Negrito apenas para cabeçalhos

def create_slide(prs, img_path, text_input, time_input, ternary_code, 
                show_chart=False, special_message=None, 
                analysis_texts=None, show_analysis=False,
                prototipo_msg=None):
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    slide.shapes.add_picture(
        img_path, Inches(0), Inches(0),
        width=Inches(10), height=Inches(5.625)
    )
    
    # Cabeçalho (com negrito)
    add_textbox(slide, text_input, 8.5, 0.2, 1.5, 0.5, 14, is_header=True)  # Nome da startup
    add_textbox(slide, time_input, 8.5, 0.7, 1.5, 0.5, 12, is_header=True)   # Carimbo
    
    # Slide 1 - Gráfico
    if show_chart:
        from chart import create_questionnaire_chart
        create_questionnaire_chart(slide, ternary_code)
    
    # Slide 2 - Respostas
    elif special_message and special_message['show_special']:
        add_textbox(
            slide, special_message['message'],
            special_message['position'][0], special_message['position'][1],
            3, 0.8, 12, is_special=True
        )
        add_textbox(
            slide, special_message['messagePadrao'],
            special_message['positionPadrao'][0], special_message['positionPadrao'][1],
            3, 0.8, 12, is_special=True
        )
    
    # Slide 2 - Análises
    elif show_analysis and analysis_texts:
        positions = [(4.2, 2.3), (4.2, 3.25), (4.11, 4.11)]
        for text, pos in zip(analysis_texts, positions):
            add_textbox(slide, text, pos[0], pos[1], 3, 0.8, 12)  # Arial 12 normal
    
    # Slide 3 - Protótipo
    elif prototipo_msg and "img3.png" in img_path:
        positions = [
            (4.2, 2),    # Título
            (4.2, 2.95)  # Conteúdo
        ]

        for i, msg in enumerate(prototipo_msg):
            offset = i * 1.5
            add_textbox(
                slide, msg['titulo'],
                positions[0][0], positions[0][1] + offset,
                4.2, 0.6, 14  # Título maior mas ainda Arial
            )
            add_textbox(
                slide, msg['conteudo'],
                positions[1][0], positions[1][1] + offset,
                3, 0.8, 12    # Conteúdo Arial 12
            )