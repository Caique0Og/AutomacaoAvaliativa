
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math
import os

def create_questionnaire_chart(slide, binary_code, center_x, center_y, radius):
    """Cria um gráfico circular no estilo do questionário com linhas e posicionamento ajustado."""
    num_questions = len(binary_code)
    circle_size = Inches(0.19)  # tamanho dos círculos 
    circle_centers = []  # Lista para armazenar as coordenadas dos centros dos círculos

    # Adicionar os círculos e armazenar suas coordenadas
    for i in range(num_questions):
        # Define o posicionamento dos círculos
        if i == 0:  # Primeiro círculo no topo
            left = center_x * 1.635
            top = center_y * 0.635
        elif i == 1:  # Segundo círculo abaixo à esquerda
            left = center_x * 1.22
            top = center_y * 0.9
        elif i == 2:  # Terceiro círculo abaixo à esquerda
            left = center_x * 0.89
            top = center_y * 1.158
        elif i == 3:  # Quarto círculo abaixo à esquerda
            left = center_x * 0.57
            top = center_y * 1.3849
        elif i == 4:  # Quinto círculo abaixo do quarto
            left = center_x * 0.57
            top = center_y * 1.6475
        elif i == 5:  # Sexto círculo abaixo do quinto
            left = center_x * 0.57
            top = center_y * 1.838
        elif i == 6:  # Sétimo círculo abaixo do sexto
            left = center_x *0.57
            top = center_y * 2.028
        elif i == 7:  # Oitavo círculo abaixo do sétimo
            left = center_x * 0.57
            top = center_y * 2.218
        elif i == 8:  # Nono círculo abaixo do oitavo
            left = center_x * 0.57
            top = center_y * 2.4
        elif i == 9:  # Décimo círculo abaixo do nono
            left = center_x * 0.57
            top = center_y * 2.590
        elif i == 10:  # Décimo primeiro círculo abaixo do décimo
            left = center_x * 0.57
            top = center_y * 2.789
        elif i == 11:  # Décimo segundo círculo
            left = center_x * 1.22
            top = center_y *  1.635
        elif i == 12:  # Décimo terceiro círculo
            left = center_x * 1.22
            top = center_y *  1.842
        elif i == 13:  # Décimo quarto círculo
            left = center_x * 1.22
            top = center_y * 2.033
        elif i == 14:  # Décimo quinto círculo
            left = center_x * 1.22
            top = center_y * 2.221
        elif i == 15:  # Décimo sexto círculo
            left = center_x * 1.22
            top = center_y * 2.4
        elif i == 16:  # Décimo sétimo círculo
            left = center_x * 1.22
            top = center_y * 2.590
        elif i == 17:  # Décimo oitavo círculo
            left = center_x *  1.635
            top = center_y * 0.9
        elif i == 18:  # Décimo nono círculo
            left = center_x * 2.05
            top = center_y * 0.9
        elif i == 19:  # Vigésimo círculo
            left = center_x * 1.39
            top = center_y * 1.6475
        elif i == 20:  # Vigésimo primeiro círculo
            left = center_x * 1.39
            top = center_y * 1.838
        elif i == 21:  # Vigésimo segundo círculo
            left = center_x *1.39
            top = center_y * 2.028
        elif i == 22:  # Vigésimo terceiro círculo
            left = center_x * 1.39
            top = center_y * 2.218
        elif i == 23:  # Vigésimo quarto círculo
            left = center_x * 1.39
            top = center_y * 2.4
        elif i == 24:  # Vigésimo quinto círculo
            left = center_x * 1.39
            top = center_y * 2.590
        elif i == 25:  # Vigésimo sexto círculo
            left = center_x * 2.381
            top = center_y * 1.159
        elif i == 26:  # Vigésimo sétimo círculo
            left = center_x * 2.04
            top = center_y * 1.637
        elif i == 27:  # Vigésimo oitavo círculo
            left = center_x * 2.04
            top = center_y * 1.837
        elif i == 28:  # Vigésimo nono círculo
            left = center_x * 2.04
            top = center_y *  2.033
        elif i == 29:  # Trigésimo círculo
            left = center_x * 2.708
            top = center_y * 1.38483
        elif i == 30: 
            left = center_x * 2.708
            top = center_y * 1.6349
        elif i == 31:  
            left = center_x * 2.708
            top = center_y * 1.842
        

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, circle_size, circle_size
        )

        # Condição para determinar a cor (0 = Vermelho, 1 = Verde) 
        if binary_code[i] == '1':
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 255, 0)  # Verde (Sim)
        else:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)  # Vermelho (Não)
        # Adicionar borda preta
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)  # Preto
        line.width = Inches(0.03)
        # Salvar as coordenadas do centro do círculo
        circle_centers.append((left + circle_size / 2, top + circle_size / 2))


        # Cria a caixa de texto para o número

        text_left = left + circle_size + Inches(0.05)
        text_top = top + circle_size/2 - Inches(0.1) 
        
        text_box = slide.shapes.add_textbox(text_left, text_top, Inches(0.5), Inches(0.2))
        text_frame = text_box.text_frame
        paragraph = text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = str(i+1)  
        font = run.font
        font.size = Inches(0.15)  
        font.bold = True

        # Cria a caixa de texto para o número

    # Adicionar as linhas entre os círculos
    for i in range(num_questions):
        # Pega o ponto de início e o ponto final para cada linha
        start_x, start_y = circle_centers[i]
        end_x, end_y = circle_centers[(i + 1) % num_questions]  # Conecta o último ao primeiro

        # Adiciona a linha
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
        )
        connector.line.color.rgb = RGBColor(0, 0, 0)  # Cor da linha (preto)

# Criar uma nova apresentação
prs = Presentation()

# Adicionar um slide em branco
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Definir as dimensões da imagem de fundo para 1920x1080 (simulado dentro do slide padrão)
width_in = Inches(10)   # Largura do slide padrão
height_in = Inches(5.625)  # Altura do slide padrão

# Adicionar a imagem como plano de fundo, movendo-a 0.7 polegadas para baixo
left = Inches(0)
top = Inches(0.7)  # Move a imagem 0.7 polegadas para baixo
slide.shapes.add_picture("img1.png", left, top, width=width_in, height=height_in)

# Input do valor binário
binary_code = input("Digite o valor binário (32 caracteres): ")

# Verificar se o valor binário tem 30 caracteres
while len(binary_code) != 32:
    print("O valor binário deve ter exatamente 32 caracteres.")
    binary_code = input("Digite novamente o valor binário: ")

# Input do texto para a caixa de texto
text_input = input("Nome da Startup: ")

# Posição e raio do gráfico (ajustar conforme necessário)
center_x = Inches(3)  # Ajustar a posição horizontal
center_y = Inches(2)  # Ajustar a posição vertical
radius = Inches(1.5)  # Ajustar o raio

create_questionnaire_chart(slide, binary_code, center_x, center_y, radius)

# Adicionar uma caixa de texto com as coordenadas
left = center_x *  2.78 
top = center_y * 0.37
text_box = slide.shapes.add_textbox(
    left, top, Inches(2), Inches(0.5)
)
text_frame = text_box.text_frame
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = text_input
font = run.font
font.bold = True


time_input = input("Carimbo data hora: ")

left_time = center_x *  2.54
top_time = center_y *  0.56



time_box = slide.shapes.add_textbox(
    left_time, top_time, Inches(2), Inches(0.5)
)
time_frame = time_box.text_frame
time_paragraph = time_frame.paragraphs[0]
time_run = time_paragraph.add_run()
time_run.text = time_input
time_font = time_run.font
time_font.bold = True

slide = prs.slides.add_slide(prs.slide_layouts[6])
width_in = Inches(10)  
height_in = Inches(5.625)
left = Inches(0)
top = Inches(0.7)  
slide.shapes.add_picture("img2.png", left, top, width=width_in, height=height_in)

center_x = Inches(3) 
center_y = Inches(2) 
radius = Inches(1.5)



left = center_x *  2.78 
top = center_y * 0.37
text_box = slide.shapes.add_textbox(
    left, top, Inches(2), Inches(0.5)
)
text_frame = text_box.text_frame
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = text_input
font = run.font
font.bold = True

left_time = center_x *  2.54
top_time = center_y *  0.56


time_box = slide.shapes.add_textbox(
    left_time, top_time, Inches(2), Inches(0.5)
)
time_frame = time_box.text_frame
time_paragraph = time_frame.paragraphs[0]
time_run = time_paragraph.add_run()
time_run.text = time_input
time_font = time_run.font
time_font.bold = True


slide = prs.slides.add_slide(prs.slide_layouts[6])
width_in = Inches(10)  
height_in = Inches(5.625)
left = Inches(0)
top = Inches(0.7)  
slide.shapes.add_picture("img3.png", left, top, width=width_in, height=height_in)

center_x = Inches(3) 
center_y = Inches(2) 
radius = Inches(1.5)

left = center_x *  2.78 
top = center_y * 0.37
text_box = slide.shapes.add_textbox(
    left, top, Inches(2), Inches(0.5)
)
text_frame = text_box.text_frame
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = text_input
font = run.font
font.bold = True

left_time = center_x *  2.54
top_time = center_y *  0.56


time_box = slide.shapes.add_textbox(
    left_time, top_time, Inches(2), Inches(0.5)
)
time_frame = time_box.text_frame
time_paragraph = time_frame.paragraphs[0]
time_run = time_paragraph.add_run()
time_run.text = time_input
time_font = time_run.font
time_font.bold = True


# Salvar a apresentação
file_path = "apresentacao_com_fundo_e_grafico.pptx"
prs.save(file_path)

# Abrir a apresentação automaticamente
os.startfile(file_path)