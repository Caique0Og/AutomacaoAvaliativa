from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from constants import GREEN, RED, BLACK, CIRCLE_SIZE, LINE_WIDTH, TEXT_FONT_SIZE, CIRCLE_POSITIONS

def create_questionnaire_chart(slide, ternary_code):
    """Cria gráfico com 3 estados possíveis"""
    num_questions = len(ternary_code)
    circle_size = CIRCLE_SIZE

    for i in range(num_questions):
        # Agora usa coordenadas absolutas em polegadas
        left = Inches(CIRCLE_POSITIONS[i][0])
        top = Inches(CIRCLE_POSITIONS[i][1])

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, circle_size, circle_size
        )

        # Define a cor baseado no valor (0, 1 ou 2)
        fill = shape.fill
        fill.solid()
        if ternary_code[i] == '1':    # Verde para Sim
            fill.fore_color.rgb = RGBColor(0, 255, 0)
        elif ternary_code[i] == '0':  # Vermelho para Não
            fill.fore_color.rgb = RGBColor(255, 0, 0)
        else:                        # Branco para Sem resposta
            fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Mantém borda preta
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = LINE_WIDTH

def add_circle_number(slide, index, left, top):
    """Adiciona o número do círculo ao lado."""
    text_box = slide.shapes.add_textbox(
        left + CIRCLE_SIZE + Inches(0.05),
        top + CIRCLE_SIZE/2 - Inches(0.1),
        Inches(0.5),
        Inches(0.2)
    )
    text_frame = text_box.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run() 
    run.text = str(index + 1)
    run.font.size = TEXT_FONT_SIZE
    run.font.bold = True
